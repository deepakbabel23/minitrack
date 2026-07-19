#!/usr/bin/env python3
"""Generate an editable MiniTrack facilitator deck (.pptx), themed in the
'MiniTrack Precision' design system. Native shapes/text => editable in
PowerPoint / Google Slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "/Users/purnimababel/Documents/GitHub/minitrack/minitrack_backend/facilitator/deck.pptx"

# ---- palette (MiniTrack Precision) --------------------------------------
def C(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
BG        = C("F5F7FB"); SURFACE = C("FFFFFF"); BORDER = C("E2E8F0")
TEXT      = C("0F172A"); MUTED   = C("475569"); SUBTLE = C("94A3B8")
ACCENT    = C("4F46E5"); ACC_STR = C("4338CA"); ACC_SOFT = C("EEF2FF"); ACC_INK = C("312E81")
GOOD      = C("16A34A"); GOOD_SOFT = C("DCFCE7"); GOOD_TX = C("166534")
DANGER    = C("DC2626"); DANG_SOFT = C("FEE2E2"); DANG_TX = C("991B1B")
WARN_SOFT = C("FEF3C7"); WARN_TX = C("92400E"); INFO_SOFT = C("DBEAFE"); INFO_TX = C("1E40AF")
CODE_BG   = C("0B1120"); CODE_SURF = C("111A2E"); CODE_TX = C("E2E8F0"); CODE_DIM = C("94A3B8")
CODE_ACC  = C("A5B4FC"); CODE_GOOD = C("86EFAC"); CODE_DANG = C("FCA5A5")
WHITE     = C("FFFFFF")

FONT = "Arial"; MONO = "Consolas"
CW, CH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(CW)
prs.slide_height = Inches(CH)
BLANK = prs.slide_layouts[6]

# ---- helpers ------------------------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    return s

def rrect(s, x, y, w, h, fill, line=None, lw=1.0, rad=0.06):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = rad
    except Exception: pass
    return sp

def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=4, space_before=0, leading=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    if leading: p.line_spacing = leading
    return p

def run(p, t, size=14, color=TEXT, bold=False, font=FONT, italic=False):
    r = p.add_run(); r.text = t
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = font; f.italic = italic; f.color.rgb = color
    return r

def eyebrow(s, t, x=0.9, y=0.62):
    w = 0.14 * len(t) + 0.5
    rrect(s, x, y, min(w, 7.5), 0.4, ACC_SOFT, rad=0.5)
    tf = box(s, x+0.18, y+0.02, min(w,7.5), 0.36, MSO_ANCHOR.MIDDLE)
    run(para(tf, True), t.upper(), 11, ACCENT, True)

def title(s, t, x=0.9, y=1.15, size=30, w=11.6):
    tf = box(s, x, y, w, 1.4)
    run(para(tf, True, leading=1.02), t, size, TEXT, True)

def brand(s, dark=False):
    rrect(s, 0.9, 6.95, 0.28, 0.28, ACCENT, rad=0.35)
    tf = box(s, 1.08, 6.93, 3.5, 0.32, MSO_ANCHOR.MIDDLE)
    p = para(tf, True)
    run(p, "  MiniTrack", 11, (WHITE if dark else TEXT), True)
    run(p, "  · Live Build Deck", 10, (C("C7CDFF") if dark else SUBTLE), False)

def pagenum(s, n, dark=False):
    tf = box(s, 12.1, 6.95, 1.1, 0.32, MSO_ANCHOR.MIDDLE)
    p = para(tf, True, align=PP_ALIGN.RIGHT)
    run(p, f"{n:02d}", 11, (WHITE if dark else TEXT), True, MONO)
    run(p, f" / {TOTAL:02d}", 10, (C("C7CDFF") if dark else SUBTLE), False, MONO)

def bullets(tf, items, size=13, gap=6):
    first = True
    for it in items:
        p = para(tf, first, space_after=gap, leading=1.12); first = False
        run(p, "▪  ", 12, ACCENT, True)
        if isinstance(it, str):
            run(p, it, size, MUTED)
        else:
            for seg, bold in it: run(p, seg, size, (TEXT if bold else MUTED), bold)

def term(s, x, y, w, h, name, lines):
    """dark terminal-style prompt block. lines = list of (text, style)."""
    rrect(s, x, y, w, h, CODE_BG, line=C("24304A"), rad=0.05)
    # title bar
    for i, col in enumerate(["FF5F57","FEBC2E","28C840"]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.22+i*0.22), Inches(y+0.2), Inches(0.11), Inches(0.11))
        d.fill.solid(); d.fill.fore_color.rgb = C(col); d.line.fill.background(); d.shadow.inherit = False
    tfn = box(s, x+1.05, y+0.12, w-1.2, 0.3, MSO_ANCHOR.MIDDLE)
    run(para(tfn, True), name, 9.5, CODE_DIM, False, MONO)
    tf = box(s, x+0.28, y+0.55, w-0.5, h-0.7)
    cmap = {"dim":CODE_DIM,"normal":CODE_TX,"accent":CODE_ACC,"good":CODE_GOOD,"danger":CODE_DANG}
    first = True
    for txt, st in lines:
        p = para(tf, first, space_after=1, leading=1.05); first = False
        run(p, txt, 10.5, cmap.get(st, CODE_TX), st in ("accent","danger"), MONO)

def callout(s, x, y, w, h, kind, label, text):
    fills = {"gotcha":(DANG_SOFT,DANGER,DANG_TX),"note":(ACC_SOFT,ACCENT,ACC_INK),"good":(GOOD_SOFT,GOOD,GOOD_TX)}
    bg, ln, tx = fills[kind]
    rrect(s, x, y, w, h, bg, line=ln, lw=1.0, rad=0.08)
    icon = {"gotcha":"⚠","note":"◆","good":"✓"}[kind]
    tfi = box(s, x+0.18, y+0.12, 0.5, h-0.2, MSO_ANCHOR.TOP)
    run(para(tfi, True), icon, 15, ln, True)
    tf = box(s, x+0.7, y+0.14, w-0.9, h-0.22, MSO_ANCHOR.MIDDLE)
    p = para(tf, True, leading=1.08)
    run(p, label.upper()+"  ", 9, tx, True)
    run(p, text, 12, tx, False)

def card(s, x, y, w, h, fill=SURFACE, line=BORDER):
    return rrect(s, x, y, w, h, fill, line=line, lw=1.0, rad=0.07)

TOTAL = 16
n = 0
def nx():
    global n; n += 1; return n

# ========================================================================
# 1 — COVER
# ========================================================================
s = slide(ACCENT)
# darker panel accent
rrect(s, 0, 0, CW, CH, ACCENT, rad=0.001)
p2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(-2.2), Inches(7), Inches(7))
p2.fill.solid(); p2.fill.fore_color.rgb = ACC_STR; p2.line.fill.background(); p2.shadow.inherit=False
tf = box(s, 0.95, 0.75, 6, 0.5)
run(para(tf, True), "K21 ACADEMY · COURSE 2606: CLAUDE AI", 12, C("C7CDFF"), True)
tf = box(s, 0.9, 1.9, 9.6, 2.4)
p = para(tf, True, leading=1.0)
run(p, "Build MiniTrack's\nfrontend — ", 46, WHITE, True)
run(p, "live.", 46, C("C7CDFF"), True)
tf = box(s, 0.95, 4.35, 8.8, 1.3)
run(para(tf, True, leading=1.25), "A finished design system on one side, a real API contract on the other. In eight prompted acts, we turn them into a complete, accessible React + Vite app — and let a reviewer agent and tests prove it's right.", 15, C("E0E2FF"))
chips = ["React + Vite + TypeScript", 'Design system: "MiniTrack Precision"', "Plan → Implement → Review → Verify"]
cx = 0.95
for ch in chips:
    w = 0.098*len(ch)+0.5
    rrect(s, cx, 5.85, w, 0.42, C("5B52E8"), rad=0.5)
    tfc = box(s, cx+0.2, 5.87, w, 0.38, MSO_ANCHOR.MIDDLE)
    run(para(tfc, True), ch, 11, WHITE, True)
    cx += w + 0.2
brand(s, dark=True);
tf = box(s, 11.7, 6.95, 1.5, 0.32, MSO_ANCHOR.MIDDLE)
run(para(tf, True, align=PP_ALIGN.RIGHT), "01 / 16", 11, C("C7CDFF"), True, MONO)
n = 1  # cover counted

# ========================================================================
# 2 — WHAT WE'RE BUILDING
# ========================================================================
s = slide(); eyebrow(s, "The target"); title(s, "Five screens. Nothing invented.")
tf = box(s, 0.9, 2.02, 11.5, 0.6)
run(para(tf, True, leading=1.2), "MiniTrack is a minimalist task tracker. You connect with an API key, then view, create, edit, complete, and delete simple tasks.", 14, MUTED)
screens = [("/connect","Enter an API key (no accounts). Show/hide field, opt-in remember, inline error."),
           ("/tasks","Filter All / Active / Completed, task cards, Load more, inline complete & delete."),
           ("/tasks/:id","One task, its badges, and only the actions that exist."),
           ("/tasks/new · edit","One shared form for create and edit, with matching validation."),
           ("*","A friendly Not-Found with a way back."),
           ("The guardrail","A task is ONLY id, title, description, priority, completed. No dates, assignees, tags — or Reopen.")]
gx, gy, gw, gh = 0.9, 2.75, 3.74, 1.75
for i,(t,d) in enumerate(screens):
    col = i % 3; rowi = i // 3
    x = gx + col*(gw+0.19); y = gy + rowi*(gh+0.2)
    guard = (t == "The guardrail")
    card(s, x, y, gw, gh, fill=(ACC_SOFT if guard else SURFACE), line=(ACCENT if guard else BORDER))
    tfh = box(s, x+0.25, y+0.22, gw-0.45, 0.4)
    run(para(tfh, True), t, 14, (ACC_INK if guard else ACCENT), True, (FONT if guard else MONO))
    tfb = box(s, x+0.25, y+0.7, gw-0.45, gh-0.8)
    run(para(tfb, True, leading=1.14), d, 12, (ACC_INK if guard else MUTED))
brand(s); pagenum(s, nx())

# ========================================================================
# 3 — TWO INPUTS
# ========================================================================
s = slide(); eyebrow(s, "Where the build starts"); title(s, "Two inputs feed everything.")
card(s, 0.9, 2.15, 5.65, 4.35)
tf = box(s, 1.2, 2.4, 5.1, 0.5); run(para(tf, True), "🎨  A design system, as CSS", 16, TEXT, True)
tf = box(s, 1.2, 2.95, 5.1, 0.6); run(para(tf, True, leading=1.15), "Generated as \"MiniTrack Precision\" and handed to us finished — we consume it as-is.", 12.5, MUTED)
bullets(box(s, 1.2, 3.75, 5.1, 2.6), [
    [("DESIGN.md", True), (" — the written spec", False)],
    [("tokens.css", True), (" — colors, type, spacing, radius", False)],
    [("typography.css", True), (" — .text-headline-* utilities", False)],
    [("components.css", True), (" — .btn · .form-field · .badge · .task-card · .dialog", False)],
], size=12.5, gap=9)
card(s, 6.78, 2.15, 5.65, 4.35)
tf = box(s, 7.08, 2.4, 5.1, 0.5); run(para(tf, True), "🔌  A real backend contract", 16, TEXT, True)
tf = box(s, 7.08, 2.95, 5.1, 0.6); run(para(tf, True, leading=1.15), "A FastAPI + SQLite JSON API — decoupled, CORS-enabled, key-protected.", 12.5, MUTED)
bullets(box(s, 7.08, 3.75, 5.1, 2.6), [
    [("The app/ layered backend", True), (" (routes → services → data)", False)],
    [("The behavioral contract in spec.md", False)],
    [("Auth via an X-API-Key header", False)],
    [("Four contract details", True), (" that bite if ignored", False)],
], size=12.5, gap=9)
brand(s); pagenum(s, nx())

# ========================================================================
# 4 — WHERE THE DESIGN CAME FROM
# ========================================================================
s = slide(); eyebrow(s, "Optional aside · design-to-code"); title(s, "The design system didn't appear by magic.", size=27)
tf = box(s, 0.9, 1.95, 11.5, 0.7)
run(para(tf, True, leading=1.18), "It was generated from a one-paragraph brief by a design tool — Google Stitch — and you can prototype the same brief in Figma Make. Same words, two tools; the tokens they produce are exactly what our CSS encodes.", 13.5, MUTED)
term(s, 0.9, 2.9, 6.55, 3.5, "Stitch / Figma Make — brief", [
    ('Design "MiniTrack" — a minimalist task', "normal"),
    ('tracker. Corporate Minimalism: near-white', "normal"),
    ('bg, white surfaces, Indigo #4f46e5 primary,', "normal"),
    ('Inter, 4px spacing, 8px cards, pill badges', "normal"),
    ('(icon + label), 3 soft elevation levels,', "normal"),
    ('centered 880px.', "normal"),
    ('# then: generate the 5 screens, no new fields', "dim"),
])
callout(s, 7.7, 2.9, 4.73, 1.55, "note", "The point", "The prototype is a picture to aim at. The CSS in src/styles/ (see preview.html) is what Acts 1–8 consume.")
callout(s, 7.7, 4.65, 4.73, 1.75, "gotcha", "Gotcha", "These tools invent due dates / assignees if asked — resist it. Never paste a real API key into a design tool.")
brand(s); pagenum(s, nx())

# ========================================================================
# 5 — PIRV
# ========================================================================
s = slide(); eyebrow(s, "The method"); title(s, "Every act is one PIRV loop.")
tf = box(s, 0.9, 2.02, 11.4, 0.6)
run(para(tf, True, leading=1.18), "We don't make the smallest possible edit. We plan the slice, implement it in the right place, review it against the contract, and verify it runs.", 14, MUTED)
steps = [("P","Plan","Name the slice and the contract rules it must honor — before any code."),
         ("I","Implement","Prompt Claude to build it, with the gotchas baked into the prompt."),
         ("R","Review","A read-only frontend-reviewer agent checks it against spec.md."),
         ("V","Verify","Run it: dev server, typecheck, lint, and the Vitest suite.")]
sx, sw = 0.9, 2.79
for i,(L,t,d) in enumerate(steps):
    x = sx + i*(sw+0.19)
    card(s, x, 2.85, sw, 2.7)
    tf = box(s, x+0.28, 3.05, 1, 0.7); run(para(tf, True), L, 34, ACCENT, True, MONO)
    tf = box(s, x+0.28, 3.85, sw-0.5, 0.4); run(para(tf, True), t, 15, TEXT, True)
    tf = box(s, x+0.28, 4.3, sw-0.5, 1.1); run(para(tf, True, leading=1.14), d, 11.5, MUTED)
    if i < 3:
        tfa = box(s, x+sw-0.02, 3.85, 0.25, 0.4); run(para(tfa, True), "→", 16, SUBTLE, True)
callout(s, 0.9, 5.8, 11.53, 0.85, "good", "Why it matters here", "The first build hit the contract's sharp edges the hard way. This time we front-load them, so the live build is right the first time.")
brand(s); pagenum(s, nx())

# ========================================================================
# 6 — DESIGN TOUR
# ========================================================================
s = slide(); eyebrow(s, "The design system, on screen"); title(s, '"MiniTrack Precision" — corporate minimalism.', size=26)
sw_items = [("Primary","4F46E5",WHITE),("Hover","4338CA",WHITE),("Success","16A34A",WHITE),
            ("Danger","DC2626",WHITE),("Surface","F8FAFC",TEXT),("Text","0F172A",WHITE)]
sx = 0.9
for nm,hx,tc in sw_items:
    rrect(s, sx, 2.05, 1.82, 1.0, C(hx), line=(BORDER if hx=="F8FAFC" else None), rad=0.09)
    tf = box(s, sx+0.12, 2.62, 1.6, 0.4, MSO_ANCHOR.MIDDLE)
    run(para(tf, True), nm, 10.5, tc, True);
    tf2 = box(s, sx, 3.1, 1.82, 0.3); run(para(tf2, True), "#"+hx.lower(), 9, SUBTLE, False, MONO)
    sx += 1.99
# buttons row
card(s, 0.9, 3.75, 6.5, 1.25)
btns = [("Create task",ACCENT,WHITE,None),("Edit",SURFACE,TEXT,BORDER),("Delete",DANGER,WHITE,None),("Cancel",BG,MUTED,None)]
bx = 1.15
for t,f,tc,ln in btns:
    w = 0.11*len(t)+0.55
    rrect(s, bx, 4.05, w, 0.5, f, line=ln, rad=0.16)
    tf = box(s, bx, 4.06, w, 0.48, MSO_ANCHOR.MIDDLE); run(para(tf, True, align=PP_ALIGN.CENTER), t, 11.5, tc, True)
    bx += w + 0.2
badges = [("▲ High","FEE2E2","991B1B"),("● Medium","FEF3C7","92400E"),("▼ Low","DBEAFE","1E40AF"),("Active","F1F5F9","475569"),("✓ Completed","DCFCE7","166534")]
bx = 1.15
for t,bgc,tc in badges:
    w = 0.1*len(t)+0.42
    rrect(s, bx, 4.55, w, 0.36, C(bgc), rad=0.5)
    tf = box(s, bx, 4.55, w, 0.36, MSO_ANCHOR.MIDDLE); run(para(tf, True, align=PP_ALIGN.CENTER), t, 10, C(tc), True)
    bx += w + 0.14
# task card mock
card(s, 7.65, 3.75, 4.78, 1.98, fill=SURFACE)
rrect(s, 7.95, 4.05, 0.32, 0.32, SURFACE, line=C("C7C4D8"), lw=1.5, rad=0.2)
tf = box(s, 8.45, 4.0, 3.8, 0.4); run(para(tf, True), "Draft the Q3 report", 13, TEXT, True)
tf = box(s, 8.45, 4.4, 3.8, 0.7); run(para(tf, True, leading=1.12), "Pull the numbers, summarize wins, flag the two risks from last review.", 11, MUTED)
rrect(s, 8.45, 5.15, 1.0, 0.32, C("FEE2E2"), rad=0.5); tf=box(s,8.45,5.15,1.0,0.32,MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),"▲ High",9.5,C("991B1B"),True)
rrect(s, 9.55, 5.15, 1.0, 0.32, C("F1F5F9"), rad=0.5); tf=box(s,9.55,5.15,1.0,0.32,MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),"Active",9.5,C("475569"),True)
tf = box(s, 0.9, 5.95, 11.5, 0.6)
run(para(tf, True, leading=1.15), "4px spacing base · soft radii (4 / 8 / 12px) · three elevation levels · badges are never color alone — always icon + label.", 12, SUBTLE)
brand(s); pagenum(s, nx())

# ========================================================================
# 7 — CONTRACT GOTCHAS
# ========================================================================
s = slide(); eyebrow(s, "The learnings — front-load these"); title(s, "Four details that bit the first build.")
gts = [("① 422 body","detail is a plain STRING, not an array. Never index detail[0].msg."),
       ("② PATCH /tasks/{id}","A full replacement — send the whole object; omitted fields reset; can't change completed."),
       ("③ 401 body","Has NO request_id (404/422 do). Treat it as optional."),
       ("④ Connect","Validate the key against GET /tasks, NOT /health (health ignores the key).")]
gx = 0.9
for i,(h,d) in enumerate(gts):
    col = i%2; rowi = i//2
    x = gx + col*(5.85); y = 2.05 + rowi*1.15
    rrect(s, x, y, 5.6, 1.0, DANG_SOFT, line=DANGER, rad=0.08)
    tf = box(s, x+0.25, y+0.13, 5.1, 0.35); run(para(tf, True), h, 12.5, DANG_TX, True)
    tf = box(s, x+0.25, y+0.5, 5.1, 0.45); run(para(tf, True, leading=1.08), d, 11, DANG_TX)
# endpoint table
rows = [("GET","/tasks","filter ?completed, page ?limit&offset — no total count"),
        ("POST","/tasks","→ 201"),
        ("PATCH","/tasks/{id}","full replacement"),
        ("POST","/tasks/{id}/complete","one-way — no reopen endpoint"),
        ("DELETE","/tasks/{id}","→ 204, needs a confirm dialog")]
tbl = s.shapes.add_table(len(rows)+1, 3, Inches(0.9), Inches(4.5), Inches(11.53), Inches(2.0)).table
tbl.columns[0].width = Inches(1.4); tbl.columns[1].width = Inches(3.1); tbl.columns[2].width = Inches(7.03)
hdr = ["METHOD","PATH","NOTE"]
for c,htxt in enumerate(hdr):
    cell = tbl.cell(0,c); cell.fill.solid(); cell.fill.fore_color.rgb = ACC_SOFT
    cell.margin_left=Pt(8); cell.margin_top=Pt(3); cell.margin_bottom=Pt(3)
    p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt
    r.font.size=Pt(10); r.font.bold=True; r.font.name=FONT; r.font.color.rgb=SUBTLE
for ri,(m,path,note) in enumerate(rows, start=1):
    for c,val in enumerate([m,path,note]):
        cell = tbl.cell(ri,c); cell.fill.solid(); cell.fill.fore_color.rgb = SURFACE
        cell.margin_left=Pt(8); cell.margin_top=Pt(2); cell.margin_bottom=Pt(2)
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
        r.font.size=Pt(11); r.font.name=(MONO if c<2 else FONT)
        r.font.bold = (c==0); r.font.color.rgb = (ACCENT if c==0 else (ACC_INK if c==1 else MUTED))
brand(s); pagenum(s, nx())

# ========================================================================
# 8 — THE ARC
# ========================================================================
s = slide(); eyebrow(s, "The build arc"); title(s, "Eight acts — one commit each.")
acts = [("ACT 1","Scaffold + Connect","ab8d9c6"),("ACT 2","Task list","9d01446"),
        ("ACT 3","Create form","1eb8c15"),("ACT 4","Task detail","2aaf787"),
        ("ACT 5","Edit + useFlash","7a30a7d"),("ACT 6","Mid-session 401","d313a63"),
        ("ACT 7","Review pass","1a0bf22"),("ACT 8","Tests + verify","13ec334")]
ax, ay, aw, ah = 0.9, 2.2, 2.79, 1.6
for i,(nn,t,sha) in enumerate(acts):
    col=i%4; rowi=i//4
    x = ax+col*(aw+0.19); y = ay+rowi*(ah+0.25)
    card(s, x, y, aw, ah)
    rrect(s, x, y, 0.08, ah, ACCENT, rad=0.0)
    tf=box(s,x+0.3,y+0.22,aw-0.5,0.3); run(para(tf,True),nn,10.5,ACCENT,True,MONO)
    tf=box(s,x+0.3,y+0.6,aw-0.5,0.55); run(para(tf,True,leading=1.05),t,14,TEXT,True)
    tf=box(s,x+0.3,y+1.15,aw-0.5,0.3); run(para(tf,True),sha,10,SUBTLE,False,MONO)
callout(s, 0.9, 5.85, 11.53, 0.85, "note", "Rhythm per act", "Say the goal → paste the prompt → watch the screen → call out the gotcha → verify. Each act is a checkpoint to stop and discuss at.")
brand(s); pagenum(s, nx())

# ========================================================================
# 9-14 — ACT SLIDES (1..6)
# ========================================================================
ACTDATA = [
 dict(nn="ACT 1", t="Scaffold + design system + Connect", sha="ab8d9c6",
   lede="Stand up Vite + React + TS in frontend/, wire in the delivered CSS, and build the API-key entry point. You connect — you don't sign in.",
   name="claude — minitrack_backend",
   lines=[("# paste into Claude Code","dim"),("Scaffold a Vite React+TS app in frontend/.","normal"),
          ("Copy DESIGN.md + src/styles/*.css in and","normal"),("import index.css once. Style with the","normal"),
          ("design system's global classes — no new","normal"),("CSS Modules. Add router + ProtectedRoute,","normal"),
          ("a fetch client (X-API-Key, status:0 on","normal"),("network fail), ApiKeyContext, and a","normal"),
          ("/connect page validating via GET /tasks.","accent")],
   bullets=[[("Wire the design system:",True),(" one index.css import",False)],
            [("Typed client:",True),(" no React import — unit-testable",False)],
            [("Key safety:",True),(" in memory; sessionStorage opt-in only",False)]],
   gotcha="Validate against GET /tasks, not /health — health accepts any key.",
   verify="npm run dev → the Connect screen renders in Indigo; a wrong key is rejected, demo-key-123 lands on /tasks."),
 dict(nn="ACT 2", t="Task list — filters & Load more", sha="9d01446",
   lede="The primary screen: URL-driven filters and pagination with no total count.",
   name="claude — frontend/",
   lines=[("Build /tasks with a .task-card list and","normal"),(".badge priority/status (icon + label,","normal"),
          ("never color alone). Filter All/Active/","normal"),("Completed in the URL (?status=).","normal"),
          ("Pagination is Load more: append, raise","accent"),("offset, hide when a short page returns,","accent"),
          ("reset on filter change. Complete in place;","normal"),("Delete via a .dialog confirm. loading /","normal"),
          ("empty / error states, aria-live.","normal")],
   bullets=[[("Filter state",True),(" lives in ?status=active",False)],
            [("Complete",True),(" → POST /tasks/{id}/complete",False)],
            [("Delete",True),(" → accessible <dialog> → 204",False)]],
   gotcha="No total → Load more, hidden on a short page, reset on filter change. No numbered pages.",
   verify="filter chips change the URL & list; Complete moves a task; Delete asks first."),
 dict(nn="ACT 3", t="Create — shared form & validation", sha="1eb8c15",
   lede="A reusable <TaskForm> and validation that matches the backend exactly.",
   name="claude — frontend/",
   lines=[("Build /tasks/new → POST /tasks (201).","normal"),("Extract a reusable <TaskForm> + <FormField>","normal"),
          ("on the .form-field classes, shared with","normal"),("edit. Title required (trim; block","normal"),
          ("whitespace-only before any request).","normal"),("Priority default medium. Disable submit","normal"),
          ("in flight. On 422, render the detail","normal"),("STRING inline — never detail[0].msg.","accent")],
   bullets=[[("Trim & block",True),(" whitespace-only titles client-side",False)],
            [("No double POST",True),(" — disable while submitting",False)],
            [("Reuse",True),(" the same form for edit",False)]],
   gotcha="422 detail is a string. detail[0].msg throws.",
   verify="a blank title is blocked before any call; a forced 422 renders its string inline."),
 dict(nn="ACT 4", t="Task detail — only real actions", sha="2aaf787",
   lede="Model the screen as explicit states, and expose only capabilities the API has.",
   name="claude — frontend/",
   lines=[("Build /tasks/:taskId → GET /tasks/{id} as","normal"),("a discriminated union: loading | loaded |","normal"),
          ("notfound | error. Show id, title,","normal"),("description, badges. Actions: Edit,","normal"),
          ("Complete (only when active), Delete","normal"),("(confirm), Back. No \"Reopen\" action.","danger"),
          ("Unknown id → 404 → \"Task not found\".","normal")],
   bullets=[[("State machine",True),(", not scattered booleans",False)],
            [("Complete",True),(" hides once done",False)],
            [("404",True),(" is a first-class state",False)]],
   gotcha="No Reopen — completion is one-way; the backend has no un-complete endpoint.",
   verify="a completed task shows no Reopen; /tasks/999999 → \"Task not found.\""),
 dict(nn="ACT 5", t="Edit — the full-replacement PATCH", sha="7a30a7d",
   lede="The subtlest rule of the day, plus a small shared-hook refactor.",
   name="claude — frontend/",
   lines=[("Build /tasks/:taskId/edit → PATCH","normal"),("/tasks/{id}, reusing <TaskForm> in edit","normal"),
          ("mode (load, pre-fill). PATCH is a FULL","accent"),("replacement: submit the entire TaskInput","accent"),
          ("every time — never partial — and never","normal"),("send or touch completed.","danger"),
          ("Extract a shared useFlash hook; reuse it.","normal")],
   bullets=[[("Load first,",True),(" then pre-fill the fields",False)],
            [("Whole object",True),(" on every save",False)],
            [("useFlash",True),(" shared across screens",False)]],
   gotcha="Send only the changed field and PATCH silently wipes the rest. Never include completed.",
   verify="edit only the title → description & priority are preserved, not reset."),
 dict(nn="ACT 6", t="Mid-session 401 — reset the session", sha="d313a63",
   lede="Handle a key that stops working after you connected — revoked or expired.",
   name="claude — frontend/",
   lines=[("Add setUnauthorizedHandler to the client.","normal"),("On any /tasks* 401 with a stored key:","accent"),
          ("clear the key (memory + sessionStorage),","normal"),("record a reason, and redirect to /connect,","normal"),
          ("which shows that reason. Remember a 401","normal"),("has no request_id — don't rely on it here.","danger")],
   bullets=[[("Global handler",True),(" in the client, not per-screen",False)],
            [("Clear + redirect,",True),(" don't just show an error",False)],
            [("Reason",True),(" surfaced on the Connect page",False)]],
   gotcha="Show an in-place error and the user is stuck on a screen that can't load. Reset the session.",
   verify="change the backend key mid-session → the next action bounces to /connect with the reason."),
]
for a in ACTDATA:
    s = slide()
    tf = box(s, 0.9, 0.62, 11.5, 0.5)
    p = para(tf, True)
    run(p, a["nn"]+"   ", 15, ACCENT, True, MONO); run(p, a["t"], 21, TEXT, True)
    tfr = box(s, 9.5, 0.66, 2.9, 0.4); run(para(tfr, True, align=PP_ALIGN.RIGHT), "↔ "+a["sha"], 12, SUBTLE, False, MONO)
    tf = box(s, 0.9, 1.35, 11.5, 0.65); run(para(tf, True, leading=1.15), a["lede"], 13.5, MUTED)
    term(s, 0.9, 2.25, 6.75, 3.55, a["name"], a["lines"])
    bullets(box(s, 8.0, 2.35, 4.4, 2.0), a["bullets"], size=12, gap=8)
    callout(s, 8.0, 4.3, 4.43, 1.5, "gotcha", "Gotcha", a["gotcha"])
    tf = box(s, 0.9, 6.05, 11.5, 0.6); run(para(tf, True, leading=1.1), "✔ Verify — "+a["verify"], 11.5, GOOD_TX)
    brand(s); pagenum(s, nx())

# ========================================================================
# 15 — ACT 7 + 8
# ========================================================================
s = slide(); eyebrow(s, "Acts 7 & 8 · Review + Verify"); title(s, "Close the loop — an agent, then tests.")
card(s, 0.9, 2.15, 5.65, 4.2)
tf=box(s,1.2,2.4,5.1,0.4); p=para(tf,True); run(p,"ACT 7 · Review   ",14,ACCENT,True); run(p,"1a0bf22",10,SUBTLE,False,MONO)
tf=box(s,1.2,2.9,5.1,0.9); run(para(tf,True,leading=1.15),"A read-only frontend-reviewer subagent checks the build against spec.md: invented features, wrong endpoints, partial PATCH, key exposure, missing states, accessibility.",12,MUTED)
bullets(box(s,1.2,3.95,5.1,2.2),[ "Move focus to \"Tasks\" + announce after actions","Extract a shared toDisplayError helper","Confirm .env is gitignored"],size=12,gap=9)
card(s, 6.78, 2.15, 5.65, 4.2)
tf=box(s,7.08,2.4,5.1,0.4); p=para(tf,True); run(p,"ACT 8 · Verify   ",14,ACCENT,True); run(p,"13ec334",10,SUBTLE,False,MONO)
tf=box(s,7.08,2.9,5.1,0.9); run(para(tf,True,leading=1.15),"A Vitest + RTL suite locks the load-bearing behaviors, plus a MANUAL_TESTING.md for the keyboard / screen-reader cases.",12,MUTED)
term(s, 7.08, 4.0, 5.05, 1.95, "frontend/", [
    ("npm run typecheck && npm run lint","normal"),("  && npm test","normal"),
    ("✓ client · connect · load-more ·","good"),("  detail · edit","good")])
callout(s, 0.9, 6.5, 11.53, 0.75, "gotcha", "Gotcha", "Native <dialog> needs a jsdom polyfill (showModal / close) or the dialog tests fail.")
brand(s); pagenum(s, nx())

# ========================================================================
# 16 — WRAP
# ========================================================================
s = slide(); eyebrow(s, "Wrap"); title(s, "What they take away.")
wrap = [("Not typing speed","The value was front-loading the contract so the build was right the first time."),
        ("The loop, not the code","Plan → Implement → Review (an agent) → Verify (tests) makes correctness repeatable."),
        ("Design as input","A delivered design system becomes a UI by wiring components to its classes.")]
wx, ww = 0.9, 3.74
for i,(t,d) in enumerate(wrap):
    x = wx + i*(ww+0.19)
    card(s, x, 2.15, ww, 2.2)
    tf=box(s,x+0.25,2.4,ww-0.45,0.6); run(para(tf,True,leading=1.05),t,15,TEXT,True)
    tf=box(s,x+0.25,3.1,ww-0.45,1.1); run(para(tf,True,leading=1.15),d,12.5,MUTED)
callout(s, 0.9, 4.65, 11.53, 0.95, "note", "To extend", "Add a new field end-to-end — backend spec.md first, then type → client → form → tests — and watch the reviewer keep you honest.")
tf = box(s, 0.9, 5.95, 11.5, 0.6)
p = para(tf, True)
run(p, "Full copy-paste prompts & the contract cheat-sheet → ", 12.5, MUTED)
run(p, "facilitator/RUNBOOK.md", 12.5, ACCENT, True, MONO)
brand(s); pagenum(s, nx())

prs.save(OUT)
print("saved", OUT, "· slides:", len(prs.slides._sldIdLst))
